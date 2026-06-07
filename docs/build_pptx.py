"""Convert docs/PROJECT_OVERVIEW.md (Marp markdown) into an EDITABLE PowerPoint.

Browser-free (unlike Marp's PPTX export, which rasterises each slide). Produces
real text boxes, real tables, and embedded PNG charts — so the deck is editable.

Run:
    python docs/build_pptx.py
Output:
    docs/PROJECT_OVERVIEW.pptx
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

DOCS = Path(__file__).resolve().parent
MD = DOCS / "PROJECT_OVERVIEW.md"
OUT = DOCS / "PROJECT_OVERVIEW.pptx"

# 16:9
SLIDE_W, SLIDE_H = Emu(12192000), Emu(6858000)
MARGIN = Emu(530000)
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xEF, 0xF6, 0xFF)


def split_slides(md: str) -> list[str]:
    # Drop YAML front matter.
    if md.startswith("---"):
        md = md.split("---", 2)[2]
    # Remove HTML comments (speaker notes).
    md = re.sub(r"<!--.*?-->", "", md, flags=re.DOTALL)
    parts = re.split(r"(?m)^---\s*$", md)
    return [p.strip() for p in parts if p.strip()]


def inline(text: str) -> str:
    # Strip markdown emphasis / code / links → plain text (keep the visible text).
    text = re.sub(r"!\[.*?\]\(.*?\)", "", text)
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = text.replace("**", "").replace("`", "")
    text = re.sub(r"(?<!\*)\*(?!\*)", "", text)
    return text.strip()


class Block:
    def __init__(self, kind, data):
        self.kind, self.data = kind, data


def parse_blocks(body: str):
    """Return (title, subtitle, [blocks]). Blocks: bullet/para/table/image/code."""
    lines = body.splitlines()
    title = subtitle = None
    blocks: list[Block] = []
    i = 0
    # Title = first heading.
    while i < len(lines) and not lines[i].strip():
        i += 1
    if i < len(lines) and lines[i].startswith("#"):
        title = inline(lines[i].lstrip("#").strip())
        i += 1
        # optional immediate ## subtitle
        while i < len(lines) and not lines[i].strip():
            i += 1
        if i < len(lines) and lines[i].startswith("##"):
            subtitle = inline(lines[i].lstrip("#").strip())
            i += 1

    buf: list[str] = []

    def flush_para():
        nonlocal buf
        if buf:
            blocks.append(Block("para", " ".join(inline(b) for b in buf).strip()))
            buf = []

    while i < len(lines):
        ln = lines[i]
        s = ln.strip()
        if s.startswith("```"):
            flush_para()
            code = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            blocks.append(Block("code", "\n".join(code)))
            continue
        m_img = re.match(r"!\[.*?\]\((.*?)\)", s)
        if m_img:
            flush_para()
            blocks.append(Block("image", m_img.group(1)))
            i += 1
            continue
        if s.startswith("|") and "|" in s[1:]:
            flush_para()
            tbl = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                row = lines[i].strip().strip("|")
                if not re.match(r"^[\s:|\-]+$", row):  # skip separator row
                    tbl.append([inline(c.strip()) for c in row.split("|")])
                i += 1
            blocks.append(Block("table", tbl))
            continue
        if re.match(r"^\s*[-*]\s+", ln) or re.match(r"^\s*\d+\.\s+", ln):
            flush_para()
            indent = (len(ln) - len(ln.lstrip())) // 2
            txt = inline(re.sub(r"^\s*([-*]|\d+\.)\s+", "", ln))
            blocks.append(Block("bullet", (indent, txt)))
            i += 1
            continue
        if s.startswith(">"):
            flush_para()
            blocks.append(Block("quote", inline(s.lstrip(">").strip())))
            i += 1
            continue
        if not s:
            flush_para()
            i += 1
            continue
        buf.append(s)
        i += 1
    flush_para()
    return title, subtitle, blocks


def add_title(slide, title, subtitle, is_section):
    tb = slide.shapes.add_textbox(MARGIN, Emu(360000), SLIDE_W - 2 * MARGIN, Emu(900000))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title or ""
    r.font.size = Pt(34 if not is_section else 40)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(22); r2.font.color.rgb = BLUE


def content_area():
    return MARGIN, Emu(1450000), SLIDE_W - 2 * MARGIN, SLIDE_H - Emu(1750000)


def add_text_blocks(slide, blocks):
    x, y, w, h = content_area()
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if b.kind == "bullet":
            indent, txt = b.data
            p.level = min(indent, 4)
            run = p.add_run(); run.text = "• " + txt
            run.font.size = Pt(18); run.font.color.rgb = DARK
            p.space_after = Pt(6)
        elif b.kind == "para":
            run = p.add_run(); run.text = b.data
            run.font.size = Pt(18); run.font.color.rgb = DARK
            p.space_after = Pt(6)
        elif b.kind == "quote":
            run = p.add_run(); run.text = "“" + b.data + "”"
            run.font.size = Pt(18); run.font.italic = True; run.font.color.rgb = BLUE
            p.space_after = Pt(6)
        elif b.kind == "code":
            for cl in b.data.splitlines():
                cp = tf.add_paragraph()
                cr = cp.add_run(); cr.text = cl
                cr.font.size = Pt(14); cr.font.name = "Consolas"; cr.font.color.rgb = GREY


def add_table(slide, rows, top=Emu(1450000)):
    if not rows:
        return
    ncol = max(len(r) for r in rows)
    rows = [r + [""] * (ncol - len(r)) for r in rows]
    nrow = len(rows)
    x, _, w, _ = content_area()
    height = min(Emu(380000) * nrow, SLIDE_H - top - MARGIN)
    gfx = slide.shapes.add_table(nrow, ncol, x, top, w, height)
    tbl = gfx.table
    for c in range(ncol):
        for r in range(nrow):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Emu(60000); cell.margin_right = Emu(60000)
            cell.margin_top = Emu(20000); cell.margin_bottom = Emu(20000)
            para = cell.text_frame.paragraphs[0]
            run = para.add_run(); run.text = rows[r][c]
            run.font.size = Pt(13)
            if r == 0:
                run.font.bold = True; run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            else:
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r % 2 else RGBColor(0xFF, 0xFF, 0xFF)


def add_image(slide, rel_path, intro_blocks):
    top = Emu(1450000)
    if intro_blocks:
        add_text_blocks(slide, intro_blocks)
    img = (DOCS / rel_path)
    if not img.exists():
        return
    # center an image up to 8" wide
    max_w = Emu(8200000)
    pic = slide.shapes.add_picture(str(img), 0, 0, width=max_w)
    pic.left = int((SLIDE_W - pic.width) / 2)
    pic.top = Emu(1650000)
    if pic.top + pic.height > SLIDE_H - MARGIN:
        scale = (SLIDE_H - MARGIN - pic.top) / pic.height
        pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
        pic.left = int((SLIDE_W - pic.width) / 2)


def main():
    slides_md = split_slides(MD.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for idx, body in enumerate(slides_md):
        slide = prs.slides.add_slide(blank)
        title, subtitle, blocks = parse_blocks(body)
        is_section = body.lstrip().startswith("# ")
        # accent bar
        bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Emu(150000))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
        add_title(slide, title, subtitle, is_section)

        imgs = [b for b in blocks if b.kind == "image"]
        tables = [b for b in blocks if b.kind == "table"]
        if imgs:
            intro = [b for b in blocks if b.kind in ("para", "bullet", "quote")]
            add_image(slide, imgs[0].data, [])
            if intro:
                # caption under image
                cap = slide.shapes.add_textbox(MARGIN, SLIDE_H - Emu(700000),
                                               SLIDE_W - 2 * MARGIN, Emu(500000))
                ctf = cap.text_frame; ctf.word_wrap = True
                cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
                cr = cp.add_run(); cr.text = " ".join(b.data if isinstance(b.data, str) else b.data[1] for b in intro)
                cr.font.size = Pt(14); cr.font.color.rgb = GREY
        elif tables:
            pre = []
            for b in blocks:
                if b.kind == "table":
                    break
                if b.kind in ("para", "bullet", "quote"):
                    pre.append(b)
            top = Emu(1450000)
            if pre:
                box = slide.shapes.add_textbox(MARGIN, Emu(1400000), SLIDE_W - 2 * MARGIN, Emu(550000))
                tf = box.text_frame; tf.word_wrap = True
                f = True
                for b in pre:
                    p = tf.paragraphs[0] if f else tf.add_paragraph(); f = False
                    txt = b.data if isinstance(b.data, str) else b.data[1]
                    run = p.add_run(); run.text = ("• " + txt) if b.kind == "bullet" else txt
                    run.font.size = Pt(15); run.font.color.rgb = DARK
                top = Emu(2050000)
            add_table(slide, tables[0].data, top=top)
            post = []
            seen_t = False
            for b in blocks:
                if b.kind == "table":
                    seen_t = True; continue
                if seen_t and b.kind == "quote":
                    post.append(b)
            if post:
                qb = slide.shapes.add_textbox(MARGIN, SLIDE_H - Emu(650000), SLIDE_W - 2 * MARGIN, Emu(480000))
                qtf = qb.text_frame; qtf.word_wrap = True
                qp = qtf.paragraphs[0]
                qr = qp.add_run(); qr.text = "“" + post[0].data + "”"
                qr.font.size = Pt(14); qr.font.italic = True; qr.font.color.rgb = BLUE
        else:
            add_text_blocks(slide, blocks)

        # footer page number
        fn = slide.shapes.add_textbox(SLIDE_W - Emu(900000), SLIDE_H - Emu(420000), Emu(700000), Emu(300000))
        fr = fn.text_frame.paragraphs[0]; fr.alignment = PP_ALIGN.RIGHT
        run = fr.add_run(); run.text = str(idx + 1)
        run.font.size = Pt(11); run.font.color.rgb = GREY

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(slides_md)} slides)")


if __name__ == "__main__":
    main()
